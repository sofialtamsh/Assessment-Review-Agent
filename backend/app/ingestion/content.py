"""Session content ingestion: extract text from .pptx / .pdf / .txt, then chunk.

Chunks carry a human-readable source reference (e.g. "Slide 4", "PDF p.3") so the
Scope agent can cite exactly where a question is (or isn't) covered.
"""
from __future__ import annotations

import hashlib
import re

from ..config import get_settings
from ..schemas import Chunk

_settings = get_settings()


def extract_segments(data: bytes, filename: str) -> list[tuple[str, str]]:
    """Return [(source_ref, text)] segments from the raw file."""
    name = (filename or "").lower()
    if name.endswith(".pptx"):
        return _from_pptx(data)
    if name.endswith(".pdf"):
        return _from_pdf(data)
    # plain text / markdown fallback
    text = data.decode("utf-8-sig", errors="replace")
    return [("Notes", text)]


def _from_pptx(data: bytes) -> list[tuple[str, str]]:
    import io

    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    segments: list[tuple[str, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for r in shape.table.rows:
                    cells = [c.text.strip() for c in r.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            segments.append((f"Slide {idx}", text))
    return segments


def _from_pdf(data: bytes) -> list[tuple[str, str]]:
    import io

    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    segments: list[tuple[str, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            segments.append((f"PDF p.{idx}", text))
    return segments


def _approx_tokens(text: str) -> int:
    # cheap proxy: ~0.75 words/token -> tokens ~= words / 0.75
    return int(len(text.split()) / 0.75) + 1


def chunk_segments(session_id: str, segments: list[tuple[str, str]]) -> list[Chunk]:
    """Split segments into ~chunk_tokens windows with overlap, preserving source refs."""
    target = _settings.embeddings.chunk_tokens
    overlap = _settings.embeddings.chunk_overlap
    chunks: list[Chunk] = []
    for source_ref, text in segments:
        sentences = _split_sentences(text)
        window: list[str] = []
        wtok = 0
        for sent in sentences:
            stok = _approx_tokens(sent)
            if wtok + stok > target and window:
                chunks.append(_mk_chunk(session_id, source_ref, " ".join(window)))
                # start next window with overlap tail
                tail, ttok = [], 0
                for s in reversed(window):
                    ttok += _approx_tokens(s)
                    tail.insert(0, s)
                    if ttok >= overlap:
                        break
                window, wtok = tail, ttok
            window.append(sent)
            wtok += stok
        if window:
            chunks.append(_mk_chunk(session_id, source_ref, " ".join(window)))
    return chunks


def _split_sentences(text: str) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    parts = re.split(r"(?<=[.!?])\s+|\n+", text)
    return [p.strip() for p in parts if p.strip()]


def _mk_chunk(session_id: str, source_ref: str, text: str) -> Chunk:
    cid = hashlib.sha1(f"{session_id}|{source_ref}|{text[:80]}".encode()).hexdigest()[:12]
    return Chunk(chunk_id=f"c_{cid}", session_id=session_id, text=text, source_ref=source_ref)


def parse_content(session_id: str, data: bytes, filename: str) -> list[Chunk]:
    return chunk_segments(session_id, extract_segments(data, filename))
