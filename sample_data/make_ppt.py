"""Generate session_ds_07.pptx — the sample session content for the demo.

Run:  python sample_data/make_ppt.py
Requires python-pptx (in backend/requirements.txt).

The "Worked Example" slide contains the exact numbers (price = 50 + 30*size,
size=4 -> 170) that assignment question q05 copies verbatim, so the Scope agent's
verbatim-lift detector has something concrete to catch.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

SLIDES = [
    ("Linear Regression — Session ds_07", [
        "Course: Machine Learning | Module: Classical ML | Unit: Regression",
        "Today: cost function, gradient descent, learning rate, R-squared, assumptions.",
    ]),
    ("The MSE Cost Function", [
        "Linear regression fits a line y = w*x + b to the data.",
        "The cost function measures how wrong the predictions are.",
        "Mean Squared Error (MSE) = average of the squared differences between predicted and actual values.",
        "Minimizing MSE finds the parameters that best fit the data.",
    ]),
    ("Gradient Descent", [
        "Gradient descent iteratively updates the parameters to reduce the cost.",
        "Each step moves the parameters in the direction that decreases MSE.",
        "As iterations proceed, the cost typically decreases toward a minimum.",
        "For linear regression the MSE cost is convex, so gradient descent reaches the global minimum.",
    ]),
    ("The Learning Rate", [
        "The learning rate (alpha) controls the size of each update step toward the minimum.",
        "If the learning rate is too high, the cost may diverge and fail to converge.",
        "If it is too low, training is slow.",
        "The learning rate is conventionally denoted by the symbol alpha.",
    ]),
    ("Worked Example: Predicting House Price", [
        "Suppose we learned the model price = 50 + 30 * size.",
        "For a house of size = 4, the predicted price = 50 + 30 * 4 = 170.",
        "This shows how the fitted line turns an input into a prediction.",
    ]),
    ("R-squared", [
        "R-squared measures the proportion of variance in the target explained by the model.",
        "An R-squared of 1.0 means the model explains all of the variance.",
        "A low R-squared means the model explains little of the variance.",
        "R-squared typically ranges between 0 and 1.",
    ]),
    ("Assumptions of Linear Regression", [
        "There is a linear relationship between the features and the target.",
        "Errors are independent and have constant variance.",
        "Understanding assumptions helps diagnose when linear regression is appropriate.",
    ]),
]


def build() -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title + Content
    for title, bullets in SLIDES:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0]
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
    out = Path(__file__).resolve().parent / "session_ds_07.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
